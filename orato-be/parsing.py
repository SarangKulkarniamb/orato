from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import pdfplumber

def normalize_bbox_pdf(x0, top, x1, bottom, page_width, page_height):
    return (
        x0 / page_width,
        top / page_height,
        (x1 - x0) / page_width,
        (bottom - top) / page_height
    )

def clean_text(text):
    if not text:
        return None
    text = text.replace("\x0b", " ").strip()
    return text if text else None

def normalize_bbox_ppt(shape, slide_width, slide_height):
    return (
        shape.left / slide_width,
        shape.top / slide_height,
        shape.width / slide_width,
        shape.height / slide_height
    )

def detect_title(candidates):
    if not candidates: return None
    candidates.sort(key=lambda x: (x["top"], -len(x["text"])))
    return candidates[0]["text"]


def parse_ppt(file_path):
    prs = Presentation(file_path)
    parsed_data = {}
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    for slide_idx, slide in enumerate(prs.slides):
        slide_id = slide_idx + 1
        parsed_data[slide_id] = {"title": None, "objects": []}
        title_candidates = []
        image_count = 0

        # --- Grab all text on the slide for context fallback ---
        slide_text_blocks = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = clean_text(shape.text)
                if t: slide_text_blocks.append(t)
        full_slide_text = " ".join(slide_text_blocks)

        for shape_idx, shape in enumerate(slide.shapes):
            obj = {
                "id": f"obj_{shape_idx}",
                "type": None,
                "text": None,
                "bbox": normalize_bbox_ppt(shape, slide_width, slide_height)
            }

            if shape.has_text_frame:
                text = clean_text(shape.text)
                if text:
                    obj["type"] = "text"
                    obj["text"] = text
                    title_candidates.append({"text": text, "top": shape.top})
                else: continue

            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                caption = None
                
                # 1. Try to find a text box directly below this image
                for other_shape in slide.shapes:
                    if other_shape.has_text_frame and other_shape.top > shape.top:
                        if other_shape.top - (shape.top + shape.height) < slide_height * 0.15: 
                            c_text = clean_text(other_shape.text)
                            if c_text: 
                                caption = f"Image explicitly showing: {c_text}"
                                break
                
                # 2. If no direct caption exists, use the whole slide's text as the image description!
                if not caption:
                    caption = f"Image diagram context: {full_slide_text[:300]}"
                
                obj["type"] = "image"
                obj["text"] = caption
                obj["image_ind"] = image_count
                image_count += 1

            elif shape.has_table:
                obj["type"] = "table"
                table_text = []
                for row in shape.table.rows:
                    row_text = [clean_text(cell.text) for cell in row.cells if clean_text(cell.text)]
                    if row_text: table_text.append(" | ".join(row_text))
                obj["text"] = "\n".join(table_text) if table_text else None
            else: continue

            parsed_data[slide_id]["objects"].append(obj)
        parsed_data[slide_id]["title"] = detect_title(title_candidates)

    return parsed_data


def parse_pdf(file_path):
    parsed_data = {}
    with pdfplumber.open(file_path) as pdf:
        for page_idx, page in enumerate(pdf.pages):
            page_id = page_idx + 1
            parsed_data[page_id] = {"title": f"Page {page_id}", "objects": []}

            # 1. Extract Text Blocks
            words = page.extract_words()
            blocks = []
            block_data = []
            
            if words:
                words.sort(key=lambda w: (round(w['top'] / 5), w['x0']))
                current_block = [words[0]]
                for word in words[1:]:
                    if word['top'] - current_block[-1]['top'] < 15:
                        current_block.append(word)
                    else:
                        blocks.append(current_block)
                        current_block = [word]
                if current_block: blocks.append(current_block)

                for i, block in enumerate(blocks):
                    text = " ".join([w["text"].strip() for w in block])
                    if len(text) < 3: continue
                    x0, top = min(w["x0"] for w in block), min(w["top"] for w in block)
                    x1, bottom = max(w["x1"] for w in block), max(w["bottom"] for w in block)
                    
                    block_data.append({"text": text, "top": top, "bottom": bottom, "x0": x0, "x1": x1})
                    bbox = normalize_bbox_pdf(x0, top, x1, bottom, page.width, page.height)

                    parsed_data[page_id]["objects"].append({
                        "id": f"block_{i}", "type": "text", "text": text, "bbox": bbox
                    })

            # --- Aggregate all page text for context fallback ---
            full_page_text = " ".join([b["text"] for b in block_data])

            # 2. Extract Images & Associate Captions
            for img_idx, img in enumerate(page.images):
                img_bottom = img['bottom']
                caption = None
                
                # Find closest text block directly below the image
                for b in block_data:
                    if b['top'] >= img_bottom and (b['top'] - img_bottom) < 100:
                        if max(img['x0'], b['x0']) < min(img['x1'], b['x1']) + 50:
                            caption = f"Image explicitly showing: {b['text']}"
                            break
                            
                # Fall back to page context
                if not caption:
                    caption = f"Image diagram context: {full_page_text[:300]}"
                            
                bbox = normalize_bbox_pdf(img['x0'], img['top'], img['x1'], img['bottom'], page.width, page.height)
                
                parsed_data[page_id]["objects"].append({
                    "id": f"img_{img_idx}",
                    "type": "image",
                    "text": caption,
                    "bbox": bbox,
                    "image_ind": img_idx
                })

    return parsed_data